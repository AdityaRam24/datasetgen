"""Binary/rich document -> plain text.

Every parser has two paths:
  1. the good library, if it happens to be installed
  2. a stdlib fallback so a bare Python install still gets *something*

`available()` reports which path is live — surfaced by `datagen doctor`.
"""

from __future__ import annotations

import html
import io
import os
import re
import zlib
from html.parser import HTMLParser
from typing import Callable

from ..util import clean_text, get_logger, zip_member_text

log = get_logger("parsers")


def _has(module: str) -> bool:
    import importlib.util

    return importlib.util.find_spec(module) is not None


def available() -> dict[str, str]:
    return {
        "pdf": "pypdf" if _has("pypdf") else "builtin (flate/text-stream)",
        "docx": "python-docx" if _has("docx") else "builtin (ooxml)",
        "pptx": "python-pptx" if _has("pptx") else "builtin (ooxml)",
        "xlsx": "openpyxl" if _has("openpyxl") else "builtin (ooxml)",
        "html": (
            "trafilatura"
            if _has("trafilatura")
            else ("beautifulsoup4" if _has("bs4") else "builtin (HTMLParser)")
        ),
    }


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def parse_pdf(data: bytes) -> str:
    if _has("pypdf"):
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            if reader.is_encrypted:
                try:
                    reader.decrypt("")  # many PDFs are "encrypted" with an empty owner pw
                except Exception:
                    log.warning("encrypted PDF could not be opened")
                    return ""
            pages = []
            for i, page in enumerate(reader.pages):
                try:
                    pages.append(page.extract_text() or "")
                except Exception as e:  # a single bad page must not kill the doc
                    log.debug("pdf page %d failed: %s", i, e)
            text = "\n\n".join(pages)
            if text.strip():
                return _postprocess_pdf(text)
        except Exception as e:
            log.warning("pypdf failed (%s) — trying builtin extractor", e)

    return _postprocess_pdf(_pdf_builtin(data))


def _pdf_builtin(data: bytes) -> str:
    """Minimal PDF text extraction: inflate every FlateDecode stream and pull
    the text-showing operators out of the content streams.

    Handles the common case (Tj / TJ / ' / ") of digitally-produced PDFs. It
    will not do OCR, CID font mapping or complex encodings — install pypdf for
    anything serious. This exists so a no-dependency run is not a dead end.
    """
    out: list[str] = []
    for match in re.finditer(rb"stream\r?\n(.*?)endstream", data, re.S):
        raw = match.group(1)
        try:
            content = zlib.decompress(raw)
        except zlib.error:
            content = raw  # uncompressed stream
        if b"Tj" not in content and b"TJ" not in content:
            continue
        out.append(_pdf_content_text(content))
    return "\n".join(t for t in out if t.strip())


_PDF_STRING = re.compile(rb"\((?:\\.|[^\\()])*\)", re.S)


def _pdf_content_text(content: bytes) -> str:
    lines: list[str] = []
    for op in re.finditer(rb"(?:\[(?P<arr>[^\]]*)\]\s*TJ)|(?:(?P<str>\((?:\\.|[^\\()])*\))\s*(?:Tj|'|\"))|(?P<nl>T\*|Td|TD|ET)", content, re.S):
        if op.group("arr") is not None:
            parts = [_pdf_decode(s) for s in _PDF_STRING.findall(op.group("arr"))]
            lines.append("".join(parts))
        elif op.group("str") is not None:
            lines.append(_pdf_decode(op.group("str")))
        else:
            lines.append("\n")
    return "".join(lines)


_PDF_ESCAPES = {b"n": "\n", b"r": "\n", b"t": "\t", b"b": "", b"f": "", b"(": "(", b")": ")", b"\\": "\\"}


def _pdf_decode(token: bytes) -> str:
    body = token[1:-1] if token.startswith(b"(") else token
    out: list[str] = []
    i = 0
    while i < len(body):
        ch = body[i : i + 1]
        if ch == b"\\" and i + 1 < len(body):
            nxt = body[i + 1 : i + 2]
            if nxt in _PDF_ESCAPES:
                out.append(_PDF_ESCAPES[nxt])
                i += 2
                continue
            if nxt.isdigit():  # \ooo octal
                octal = body[i + 1 : i + 4]
                try:
                    out.append(chr(int(octal, 8)))
                    i += 1 + len(octal)
                    continue
                except ValueError:
                    pass
            i += 2
            continue
        out.append(ch.decode("latin-1", "replace"))
        i += 1
    return "".join(out)


_PAGE_NUM = re.compile(r"^\s*(?:page\s*)?\d+\s*(?:/\s*\d+)?\s*$", re.I)


def _postprocess_pdf(text: str) -> str:
    """Repair the usual PDF damage: hyphenated line breaks, stray page numbers,
    and lines broken mid-sentence by the layout engine."""
    if not text:
        return ""
    text = text.replace("­", "")                      # soft hyphen
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)           # de-hyphenate wraps
    lines = [ln.rstrip() for ln in text.split("\n")]
    kept = [ln for ln in lines if not _PAGE_NUM.match(ln)]

    merged: list[str] = []
    for line in kept:
        if (
            merged
            and merged[-1]
            and not merged[-1].endswith((".", ":", "?", "!", ";"))
            and line
            and line[0].islower()
        ):
            merged[-1] += " " + line.strip()
        else:
            merged.append(line)
    return clean_text("\n".join(merged))


# ---------------------------------------------------------------------------
# OOXML: docx / pptx / xlsx
# ---------------------------------------------------------------------------

_XML_TAG = re.compile(r"<[^>]+>")


def _ooxml_text(data: bytes, member_pattern: str, para_tag: str) -> str:
    parts = []
    for xml in zip_member_text(data, member_pattern):
        xml = re.sub(rf"</{para_tag}>", "\n", xml)
        xml = re.sub(r"</w:tr>|</a:p>|<w:br/>", "\n", xml)
        parts.append(html.unescape(_XML_TAG.sub("", xml)))
    return clean_text("\n".join(parts))


def parse_docx(data: bytes) -> str:
    if _has("docx"):
        try:
            import docx  # python-docx

            d = docx.Document(io.BytesIO(data))
            blocks = [p.text for p in d.paragraphs]
            for table in d.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        blocks.append(" | ".join(cells))
            return clean_text("\n".join(blocks))
        except Exception as e:
            log.warning("python-docx failed (%s) — using builtin", e)
    return _ooxml_text(data, r"word/(document|header\d*|footer\d*)\.xml", "w:p")


def parse_pptx(data: bytes) -> str:
    if _has("pptx"):
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            blocks = []
            for i, slide in enumerate(prs.slides, 1):
                blocks.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        blocks.append(shape.text_frame.text)
            return clean_text("\n".join(blocks))
        except Exception as e:
            log.warning("python-pptx failed (%s) — using builtin", e)
    return _ooxml_text(data, r"ppt/slides/slide\d+\.xml", "a:p")


def parse_xlsx(data: bytes) -> str:
    if _has("openpyxl"):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            blocks = []
            for ws in wb.worksheets:
                blocks.append(f"## Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        blocks.append(" | ".join(cells))
            wb.close()
            return clean_text("\n".join(blocks))
        except Exception as e:
            log.warning("openpyxl failed (%s) — using builtin", e)
    return _ooxml_text(data, r"xl/sharedStrings\.xml", "si")


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------


class _TextExtractor(HTMLParser):
    SKIP = {"script", "style", "noscript", "svg", "head", "nav", "footer", "form"}
    BLOCK = {"p", "div", "br", "li", "tr", "section", "article", "h1", "h2", "h3", "h4", "h5", "h6"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.title = ""
        self._skip_depth = 0
        self._in_title = False
        self._heading = ""

    def handle_starttag(self, tag, attrs):
        if tag in self.SKIP:
            self._skip_depth += 1
        elif tag == "title":
            self._in_title = True
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self._heading = "#" * int(tag[1]) + " "
            self.parts.append("\n\n" + self._heading)
        elif tag in self.BLOCK:
            self.parts.append("\n")
        elif tag == "li":
            self.parts.append("\n- ")

    def handle_endtag(self, tag):
        if tag in self.SKIP and self._skip_depth:
            self._skip_depth -= 1
        elif tag == "title":
            self._in_title = False
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self._heading = ""
            self.parts.append("\n")

    def handle_data(self, data):
        if self._skip_depth:
            return
        if self._in_title:
            self.title += data
        elif data.strip():
            self.parts.append(data)

    def text(self) -> str:
        return clean_text("".join(self.parts))


def parse_html(data: bytes | str, url: str = "") -> tuple[str, str]:
    """Return (title, text). Prefers trafilatura's boilerplate removal."""
    raw = data.decode("utf-8", "replace") if isinstance(data, bytes) else data

    if _has("trafilatura"):
        try:
            import trafilatura

            extracted = trafilatura.extract(
                raw,
                url=url or None,
                include_comments=False,
                include_tables=True,
                favor_recall=True,
            )
            if extracted and len(extracted) > 200:
                meta_title = ""
                try:
                    meta = trafilatura.extract_metadata(raw)
                    meta_title = (meta.title or "") if meta else ""
                except Exception:
                    pass
                title = meta_title or _html_title(raw)
                return title, clean_text(extracted)
        except Exception as e:
            log.debug("trafilatura failed: %s", e)

    if _has("bs4"):
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(raw, "lxml" if _has("lxml") else "html.parser")
            for tag in soup(["script", "style", "noscript", "svg", "nav", "footer", "form"]):
                tag.decompose()
            title = soup.title.get_text(strip=True) if soup.title else ""
            main = soup.find("main") or soup.find("article") or soup.body or soup
            return title, clean_text(main.get_text("\n"))
        except Exception as e:
            log.debug("bs4 failed: %s", e)

    parser = _TextExtractor()
    try:
        parser.feed(raw)
    except Exception as e:
        log.debug("HTMLParser failed: %s", e)
    return parser.title.strip(), parser.text()


def _html_title(raw: str) -> str:
    m = re.search(r"<title[^>]*>(.*?)</title>", raw, re.S | re.I)
    return html.unescape(m.group(1)).strip() if m else ""


def html_links(raw: str, base_url: str) -> list[str]:
    """Absolute links from a page, for the crawler."""
    from urllib.parse import urldefrag, urljoin

    out: list[str] = []
    for m in re.finditer(r'<a\s[^>]*href=["\']([^"\'>]+)["\']', raw, re.I):
        href = html.unescape(m.group(1)).strip()
        if not href or href.startswith(("#", "javascript:", "mailto:", "tel:", "data:")):
            continue
        absolute, _ = urldefrag(urljoin(base_url, href))
        if absolute.startswith(("http://", "https://")):
            out.append(absolute)
    # De-duplicate, preserving discovery order.
    return list(dict.fromkeys(out))


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

EXTENSION_MAP = {
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".webp": "image",
    ".gif": "image",
    ".bmp": "image",
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".htm": "html",
    ".html": "html",
    ".md": "markdown",
    ".markdown": "markdown",
    ".txt": "text",
    ".log": "text",
    ".rst": "text",
    ".csv": "csv",
    ".json": "json",
    ".yaml": "text",
    ".yml": "text",
    ".sh": "code",
    ".py": "code",
    ".ts": "code",
    ".js": "code",
    ".sql": "code",
}


def parse_bytes(data: bytes, kind: str, url: str = "") -> tuple[str, str]:
    """(title, text) for a blob of a known kind."""
    if kind == "image":
        return "", parse_image(data, url)
    if kind == "pdf":
        return "", parse_pdf(data)
    if kind == "docx":
        return "", parse_docx(data)
    if kind == "pptx":
        return "", parse_pptx(data)
    if kind == "xlsx":
        return "", parse_xlsx(data)
    if kind == "html":
        return parse_html(data, url)
    if kind in ("markdown", "text"):
        text = clean_text(data.decode("utf-8", "replace"))
        return markdown_title(text), text
    if kind == "csv":
        return "", _parse_csv(data)
    if kind == "json":
        return "", _parse_json(data)
    return "", clean_text(data.decode("utf-8", "replace"))


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------

# Images are the one format that cannot be parsed — only *described*, by a local
# vision model. Parsers must stay dependency-free (the test suite runs them with
# no LLM and no network), so the describer is injected instead of imported:
# LocalLLM registers itself on construction, and every connector that calls
# parse_bytes gets image support for free — local files, crawled pages and
# Confluence attachments alike.
_IMAGE_DESCRIBER: "Callable[[bytes, str, str], str] | None" = None

IMAGE_MIME = {
    "image": "image/png", ".png": "image/png", ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg", ".webp": "image/webp", ".gif": "image/gif",
    ".bmp": "image/bmp",
}


def set_image_describer(fn: "Callable[[bytes, str, str], str] | None") -> None:
    global _IMAGE_DESCRIBER
    _IMAGE_DESCRIBER = fn


def image_support() -> bool:
    return _IMAGE_DESCRIBER is not None


def parse_image(data: bytes, url: str = "") -> str:
    """Describe an image as text. Empty string when no describer is registered,
    which makes the document fall out of the pipeline exactly like a scanned PDF
    with no text layer — nothing crashes, nothing silently ships as garbage."""
    if _IMAGE_DESCRIBER is None:
        log.debug("no vision model registered — skipping image %s", url or "<bytes>")
        return ""
    ext = os.path.splitext(url.split("?")[0])[1].lower()
    try:
        return clean_text(_IMAGE_DESCRIBER(data, url, IMAGE_MIME.get(ext, "image/png")) or "")
    except Exception as e:  # noqa: BLE001 - a vision failure must not stop a build
        log.warning("image description failed for %s: %s", url or "<bytes>", e)
        return ""


_ATX = re.compile(r"^\s{0,3}#{1,2}\s+(\S.*?)\s*#*\s*$")
_SETEXT = re.compile(r"^\s{0,3}[=-]{3,}\s*$")


def markdown_title(text: str) -> str:
    """First heading of a markdown/text document, or "".

    Without this a runbook whose first line is `# MLIS endpoint returns 503` is
    titled "mlis" after its filename — and that title is what ends up in every
    citation, every glossary source list and the Kalam knowledge base.
    """
    lines = text.split("\n", 40)[:40]
    for i, line in enumerate(lines):
        if not line.strip():
            continue
        m = _ATX.match(line)
        if m:
            return clean_text(m.group(1))[:200]
        # Setext: a line of === or --- directly underneath the title.
        if i + 1 < len(lines) and _SETEXT.match(lines[i + 1]) and len(line.strip()) <= 200:
            return clean_text(line)[:200]
        break        # a non-heading first line means there is no title to take
    return ""


def _parse_csv(data: bytes) -> str:
    import csv

    text = data.decode("utf-8", "replace")
    try:
        rows = list(csv.reader(io.StringIO(text)))
    except csv.Error:
        return clean_text(text)
    if not rows:
        return ""
    header = rows[0]
    lines = [" | ".join(header), "-" * 40]
    for row in rows[1:]:
        lines.append(" | ".join(row))
    return clean_text("\n".join(lines))


def _parse_json(data: bytes) -> str:
    import json as _json

    try:
        return clean_text(_json.dumps(_json.loads(data.decode("utf-8", "replace")), indent=2))
    except (ValueError, UnicodeDecodeError):
        return clean_text(data.decode("utf-8", "replace"))
